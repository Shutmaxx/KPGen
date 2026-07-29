# -*- coding: utf-8 -*-
"""Сборка коммерческого предложения и презентации из шаблонов."""
from __future__ import annotations

import re
from dataclasses import dataclass
from datetime import date
from pathlib import Path

from docx import Document
from pptx import Presentation
from pptx.util import Emu

from .analyze import CallAnalysis
from .dadata import Company
from .docx_utils import replace_everywhere
from .settings import Manager

# Текст-заглушка клиента в шаблоне презентации — ищем его, чтобы заменить.
PRES_CLIENT_MARKERS = ("ТРЕХГОРНАЯ", "ТРЁХГОРНАЯ", "МАНУФАКТУРА")
PRES_INN_PREFIX = "ИНН"


class DocumentGenerationError(Exception):
    """Ошибка сборки документа с понятным пользователю описанием."""


@dataclass
class GeneratedFiles:
    kp_path: Path
    presentation_path: Path
    summary_path: Path


def safe_folder_name(name: str) -> str:
    """Имя папки без символов, запрещённых в Windows."""
    cleaned = re.sub(r'[<>:"/\\|?*]', "", name)
    cleaned = re.sub(r"\s+", " ", cleaned).strip(" .")
    return (cleaned or "Компания")[:80]


def build_segment_value(company: Company, analysis: CallAnalysis) -> str:
    """Чем площадка полезна компании этого профиля.

    Регион и адрес сознательно не используются: в КП важна отраслевая
    польза, а не география клиента.
    """
    value = (analysis.segment_value or "").strip().rstrip(".")
    if value:
        return value + "."
    return (
        f"{company.industry_phrase()}, для которого закупки — "
        f"это непрерывный поток процедур."
    )


def _placeholder_map(company: Company, analysis: CallAnalysis,
                     manager: Manager) -> dict[str, str]:
    return {
        "{{ПОЛЬЗА_В_СЕГМЕНТЕ}}": build_segment_value(company, analysis),
        "{{ПОТРЕБНОСТИ_КЛИЕНТА}}": analysis.client_needs,
        "{{ЧТО_ПРЕДЛАГАЕМ}}": analysis.our_offer,
        "{{КОМПАНИЯ_КРАТКО}}": company.display_name,
        "{{КОМПАНИЯ_ПОЛНОЕ}}": company.name_full or company.display_name,
        "{{ИНН}}": company.inn,
        "{{КПП}}": company.kpp,
        "{{МЕНЕДЖЕР_ФИО}}": manager.full_name,
        "{{МЕНЕДЖЕР_ДОЛЖНОСТЬ}}": manager.position,
        "{{МЕНЕДЖЕР_ТЕЛЕФОН_1}}": manager.phone_1,
        "{{МЕНЕДЖЕР_ТЕЛЕФОН_2}}": manager.phone_2,
        "{{МЕНЕДЖЕР_ТЕЛЕФОН_3}}": manager.phone_3,
        "{{МЕНЕДЖЕР_EMAIL}}": manager.email,
        "{{МЕНЕДЖЕР_САЙТ}}": manager.site,
    }


def generate_kp(template_path: str | Path, output_path: Path, company: Company,
                analysis: CallAnalysis, manager: Manager) -> Path:
    """Собирает КП из шаблона, сохраняя бланк и оформление."""
    template = Path(template_path)
    if not template.exists():
        raise DocumentGenerationError(
            f"Не найден шаблон коммерческого предложения: {template}. "
            f"Укажите правильный путь в настройках."
        )

    try:
        document = Document(str(template))
    except Exception as exc:
        raise DocumentGenerationError(
            f"Не удалось открыть шаблон КП «{template.name}». "
            f"Возможно, файл повреждён или занят другой программой. Причина: {exc}"
        ) from exc

    replace_everywhere(document, _placeholder_map(company, analysis, manager))

    try:
        output_path.parent.mkdir(parents=True, exist_ok=True)
        document.save(str(output_path))
    except PermissionError as exc:
        raise DocumentGenerationError(
            f"Файл «{output_path.name}» открыт в Word. Закройте его и повторите."
        ) from exc
    except OSError as exc:
        raise DocumentGenerationError(
            f"Не удалось сохранить КП: {exc}"
        ) from exc

    return output_path


def _replace_slide_text(slide, company: Company) -> None:
    """Подставляет название и ИНН клиента на титульном слайде."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text

        if any(marker in text.upper() for marker in PRES_CLIENT_MARKERS):
            _set_shape_text(shape, company.display_name)
        elif text.strip().upper().startswith(PRES_INN_PREFIX):
            _set_shape_text(shape, f"ИНН: {company.inn}")


def _set_shape_text(shape, value: str) -> None:
    """Меняет текст фигуры, сохраняя начертание первого фрагмента."""
    frame = shape.text_frame
    paragraph = frame.paragraphs[0]
    if paragraph.runs:
        paragraph.runs[0].text = value
        for run in paragraph.runs[1:]:
            run.text = ""
    else:
        paragraph.add_run().text = value
    # Лишние абзацы внутри фигуры удаляем — иначе останется старый текст.
    for extra in list(frame.paragraphs[1:]):
        extra._p.getparent().remove(extra._p)


def _replace_client_logo(slide, logo_path: Path | None) -> None:
    """Заменяет логотип клиента на титульном слайде либо удаляет его.

    Логотипом считается самая крупная картинка в нижней половине слайда —
    в шаблоне это изображение клиента.
    """
    pictures = [s for s in slide.shapes if s.shape_type == 13]  # PICTURE
    if not pictures:
        return

    candidates = [p for p in pictures if p.top and p.top > Emu(1_200_000)]
    if not candidates:
        return

    target = max(candidates, key=lambda s: (s.width or 0) * (s.height or 0))
    left, top, width, height = target.left, target.top, target.width, target.height
    target._element.getparent().remove(target._element)

    if logo_path is None:
        return

    try:
        picture = slide.shapes.add_picture(str(logo_path), left, top)
    except Exception as exc:
        raise DocumentGenerationError(
            f"Не удалось вставить логотип «{logo_path.name}». "
            f"Поддерживаются PNG, JPG, GIF, BMP. Причина: {exc}"
        ) from exc

    # Вписываем в исходную рамку, сохраняя пропорции.
    scale = min(width / picture.width, height / picture.height)
    picture.width = int(picture.width * scale)
    picture.height = int(picture.height * scale)
    picture.left = left + (width - picture.width) // 2
    picture.top = top + (height - picture.height) // 2


def _replace_contacts(slide, manager: Manager) -> None:
    """Подставляет контакты менеджера на последнем слайде."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text
        if "@" not in text and "+7" not in text:
            continue

        lines = [
            manager.full_name,
            manager.position,
            manager.phone_1,
            manager.phone_2,
            manager.phone_3,
        ]
        if "@" in text and len(text.strip().splitlines()) <= 2:
            lines = [manager.email]

        frame = shape.text_frame
        first = frame.paragraphs[0]
        template_run = first.runs[0] if first.runs else None

        for extra in list(frame.paragraphs[1:]):
            extra._p.getparent().remove(extra._p)

        if template_run is None:
            first.add_run().text = lines[0]
        else:
            template_run.text = lines[0]
            for run in first.runs[1:]:
                run.text = ""

        for line in lines[1:]:
            paragraph = frame.add_paragraph()
            run = paragraph.add_run()
            run.text = line
            if template_run is not None:
                run.font.size = template_run.font.size
                run.font.name = template_run.font.name
                run.font.bold = template_run.font.bold
                if template_run.font.color and template_run.font.color.type is not None:
                    run.font.color.rgb = template_run.font.color.rgb


def _add_benefits_slide(presentation, company: Company,
                        benefits: list[str]) -> None:
    """Добавляет слайд с преимуществами под задачи клиента.

    Слайд собирается в фирменных цветах ЕСТП и ставится сразу после
    титульного, чтобы клиент увидел выгоду до общих разделов.
    """
    if not benefits:
        return

    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt

    dark = RGBColor(0x4E, 0x5C, 0x6B)
    text_color = RGBColor(0x3A, 0x3A, 0x3A)
    accent = RGBColor(0x5E, 0xAE, 0xC4)
    muted = RGBColor(0x6B, 0x72, 0x7A)

    width = presentation.slide_width
    height = presentation.slide_height

    # Берём макет из шаблона: у фирменных презентаций он обычно один,
    # поэтому обращаться к «пустому макету» по индексу нельзя.
    layouts = presentation.slide_layouts
    existing = list(presentation.slides)
    if existing:
        layout = existing[0].slide_layout
    elif len(layouts) > 6:
        layout = layouts[6]
    else:
        layout = layouts[0]

    slide = presentation.slides.add_slide(layout)
    # Плейсхолдеры макета не нужны — слайд собирается вручную.
    for shape in list(slide.placeholders):
        shape._element.getparent().remove(shape._element)

    # Заголовок
    title_box = slide.shapes.add_textbox(
        Emu(int(width * 0.06)), Emu(int(height * 0.08)),
        Emu(int(width * 0.88)), Emu(int(height * 0.16)))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_paragraph = title_frame.paragraphs[0]
    title_run = title_paragraph.add_run()
    title_run.text = "ЧТО МЫ ПРЕДЛАГАЕМ"
    title_run.font.size = Pt(28)
    title_run.font.bold = True
    title_run.font.color.rgb = dark
    title_run.font.name = "Arial"

    subtitle_box = slide.shapes.add_textbox(
        Emu(int(width * 0.06)), Emu(int(height * 0.23)),
        Emu(int(width * 0.88)), Emu(int(height * 0.09)))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    subtitle_run = subtitle_frame.paragraphs[0].add_run()
    subtitle_run.text = f"Под задачи {company.display_name}"
    subtitle_run.font.size = Pt(14)
    subtitle_run.font.color.rgb = muted
    subtitle_run.font.name = "Arial"

    # Пункты преимуществ
    top = int(height * 0.36)
    step = int(height * 0.13)
    for index, benefit in enumerate(benefits[:4]):
        marker = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Emu(int(width * 0.07)), Emu(top + index * step),
            Emu(int(width * 0.022)), Emu(int(height * 0.045)))
        marker.fill.solid()
        marker.fill.fore_color.rgb = accent
        marker.line.fill.background()
        marker.shadow.inherit = False

        text_box = slide.shapes.add_textbox(
            Emu(int(width * 0.11)), Emu(top + index * step - int(height * 0.012)),
            Emu(int(width * 0.82)), Emu(int(height * 0.09)))
        frame = text_box.text_frame
        frame.word_wrap = True
        run = frame.paragraphs[0].add_run()
        run.text = benefit
        run.font.size = Pt(16)
        run.font.color.rgb = text_color
        run.font.name = "Arial"

    # Декоративная «волна» в фирменном стиле
    wave = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Emu(int(width * 0.78)), Emu(int(height * 0.86)),
        Emu(int(width * 0.4)), Emu(int(height * 0.34)))
    wave.fill.solid()
    wave.fill.fore_color.rgb = accent
    wave.line.fill.background()
    wave.shadow.inherit = False

    brand_box = slide.shapes.add_textbox(
        Emu(int(width * 0.03)), Emu(int(height * 0.03)),
        Emu(int(width * 0.2)), Emu(int(height * 0.05)))
    brand_run = brand_box.text_frame.paragraphs[0].add_run()
    brand_run.text = "ESTP.RU"
    brand_run.font.size = Pt(10)
    brand_run.font.color.rgb = muted
    brand_run.font.name = "Arial"

    # Переносим слайд на третью позицию — сразу после титульного клиента.
    slide_ids = presentation.slides._sldIdLst
    new_slide = slide_ids[-1]
    slide_ids.remove(new_slide)
    slide_ids.insert(2, new_slide)


def generate_presentation(template_path: str | Path, output_path: Path,
                          company: Company, manager: Manager,
                          logo_path: str | Path | None = None,
                          benefits: list[str] | None = None) -> Path:
    """Собирает презентацию: титул клиента, логотип, контакты менеджера."""
    template = Path(template_path)
    if not template.exists():
        raise DocumentGenerationError(
            f"Не найден шаблон презентации: {template}. "
            f"Укажите правильный путь в настройках."
        )

    try:
        presentation = Presentation(str(template))
    except Exception as exc:
        raise DocumentGenerationError(
            f"Не удалось открыть шаблон презентации «{template.name}». Причина: {exc}"
        ) from exc

    slides = list(presentation.slides)
    if len(slides) < 2:
        raise DocumentGenerationError(
            "В шаблоне презентации меньше двух слайдов — "
            "не найден титульный слайд клиента."
        )

    logo = Path(logo_path) if logo_path else None
    if logo is not None and not logo.exists():
        raise DocumentGenerationError(f"Файл логотипа не найден: {logo}")

    client_slide = slides[1]
    _replace_slide_text(client_slide, company)
    _replace_client_logo(client_slide, logo)
    _replace_contacts(slides[-1], manager)
    _add_benefits_slide(presentation, company, benefits or [])

    try:
        output_path.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(str(output_path))
    except PermissionError as exc:
        raise DocumentGenerationError(
            f"Файл «{output_path.name}» открыт в PowerPoint. Закройте его и повторите."
        ) from exc
    except OSError as exc:
        raise DocumentGenerationError(f"Не удалось сохранить презентацию: {exc}") from exc

    return output_path


def generate_all(company: Company, analysis: CallAnalysis, manager: Manager,
                 kp_template: str, presentation_template: str, output_root: str,
                 logo_path: str | None = None,
                 transcript: str = "") -> GeneratedFiles:
    """Готовит все файлы по одному разговору в отдельной папке."""
    folder = Path(output_root) / f"{safe_folder_name(company.display_name)}_{date.today():%Y-%m-%d}"
    folder.mkdir(parents=True, exist_ok=True)

    kp_path = generate_kp(kp_template, folder / "Коммерческое предложение.docx",
                          company, analysis, manager)
    pres_path = generate_presentation(presentation_template, folder / "Презентация.pptx",
                                      company, manager, logo_path,
                                      benefits=analysis.benefits)

    summary_path = folder / "Выжимка разговора.txt"
    summary_path.write_text(analysis.summary + "\n", encoding="utf-8")

    if transcript:
        (folder / "Расшифровка разговора.txt").write_text(transcript, encoding="utf-8")

    return GeneratedFiles(kp_path, pres_path, summary_path)
