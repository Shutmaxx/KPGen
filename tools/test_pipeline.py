# -*- coding: utf-8 -*-
"""Сквозная проверка конвейера на готовой расшифровке (без GUI).

Использование:
    py tools/test_pipeline.py <ИНН> <файл расшифровки>
"""
from __future__ import annotations

import sys
import time
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))
sys.stdout.reconfigure(encoding="utf-8")

from app.services import analyze, dadata, docgen, normalize
from app.services.settings import Settings

ROOT = Path(__file__).resolve().parents[2]

inn = sys.argv[1] if len(sys.argv) > 1 else "2352002170"
transcript_file = Path(sys.argv[2]) if len(sys.argv) > 2 \
    else ROOT / "transcript_fanagoria.txt"

settings = Settings()
settings.output_dir = str(Path(__file__).resolve().parents[1] / "_test_output")

print("1) Поиск компании по ИНН…")
company = dadata.find_by_inn(inn, settings.dadata_token)
print("   ", company.display_name, "|", company.industry_phrase())

print("2) Нормализация расшифровки…")
raw = transcript_file.read_text(encoding="utf-8")
transcript = normalize.normalize(
    normalize.strip_timecodes(raw),
    normalize.company_aliases(company.display_name),
)

print("3) Анализ разговора…")
t0 = time.time()
result = analyze.analyze_call(
    transcript, company,
    use_ai=settings.use_ai,
    base_url=settings.ollama_url,
    model=settings.ollama_model,
    progress=lambda p, m: print(f"     {p:3d}% {m}"),
)
print(f"   готово за {time.time() - t0:.0f} c | ИИ: {result.used_ai}")
print()
print("   ВЫЖИМКА:", result.summary)
print("   предложений:", len(analyze.split_sentences(result.summary)))
print()
print("   ПОЛЬЗА В СЕГМЕНТЕ:", result.segment_value)
print("   ПОТРЕБНОСТИ:", result.client_needs)
print("   ПРЕДЛОЖЕНИЕ:", result.our_offer)
print("   ПРЕИМУЩЕСТВА:", result.benefits)
if result.notes:
    print("   замечания:", result.notes)

print("\n4) Сборка документов…")
files = docgen.generate_all(
    company, result, settings.manager,
    settings.kp_template, settings.presentation_template,
    settings.output_dir, logo_path=None, transcript=transcript,
)
print("   КП         :", files.kp_path)
print("   Презентация:", files.presentation_path)

print("\n5) Контроль КП…")
from docx import Document
doc = Document(str(files.kp_path))
left = [p.text for p in doc.paragraphs if "{{" in p.text]
print("   незаполненных плейсхолдеров:", len(left), left[:3])
for index in (2, 3, 10, 18):
    paragraph = doc.paragraphs[index]
    bold_map = [(r.text[:32], r.bold) for r in paragraph.runs if r.text.strip()]
    print(f"   [{index}] {paragraph.text[:105]}")
    if index == 2:
        print(f"        начертание: {bold_map}")

print("\n6) Контроль презентации…")
from pptx import Presentation
pres = Presentation(str(files.presentation_path))
slides = list(pres.slides)
print("   слайдов:", len(slides))
for number in (1, 2):
    texts = [s.text_frame.text.strip().replace("\n", " / ")
             for s in slides[number].shapes
             if s.has_text_frame and s.text_frame.text.strip()]
    print(f"   слайд {number + 1}: {texts}")
